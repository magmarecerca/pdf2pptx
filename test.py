import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image
import os


def pdf_to_pptx(pdf_path, pptx_path, dpi=300):
    print(f"Starting conversion: {pdf_path}")

    doc = fitz.open(pdf_path)

    presentation = Presentation()

    first_page = doc[0]
    pix = first_page.get_pixmap(dpi=dpi)
    img = Image.open(io.BytesIO(pix.tobytes()))
    img_width, img_height = img.size

    aspect_ratio = img_width / img_height

    MAX_SLIDE_SIZE = 56
    slide_width_inches = img_width / dpi * 2.54
    slide_height_inches = img_height / dpi * 2.54

    if slide_width_inches > MAX_SLIDE_SIZE:
        slide_width_inches = MAX_SLIDE_SIZE
        slide_height_inches = slide_width_inches / aspect_ratio

    if slide_height_inches > MAX_SLIDE_SIZE:
        slide_height_inches = MAX_SLIDE_SIZE
        slide_width_inches = slide_height_inches * aspect_ratio

    presentation.slide_width = Inches(slide_width_inches)
    presentation.slide_height = Inches(slide_height_inches)

    iteration = 0
    for page in doc:
        iteration += 1
        print(f'Processing page {iteration}/{len(doc)} of {pdf_path}')

        pix = page.get_pixmap(dpi=dpi)
        img = Image.open(io.BytesIO(pix.tobytes()))

        temp_image_path = "temp_slide.png"
        img.save(temp_image_path, "PNG")

        slide = presentation.slides.add_slide(presentation.slide_layouts[5])

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.text_frame.text == "":
                slide.shapes._spTree.remove(shape._element)

        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        img_width, img_height = img.size

        scale = min(slide_width / img_width, slide_height / img_height)
        new_width = img_width * scale
        new_height = img_height * scale

        left = (slide_width - new_width) / 2
        top = (slide_height - new_height) / 2

        slide.shapes.add_picture(temp_image_path, left, top, new_width, new_height)

    presentation.save(pptx_path)
    print(f"Finished conversion: {pdf_path} -> {pptx_path}")


def convert_pdfs_in_directory(root_dir, dpi=300):
    pdf_files = []
    for folder, _, files in os.walk(root_dir):
        for file in files:
            if file.lower().endswith(".pdf"):
                pdf_path = os.path.join(folder, file)
                pptx_path = os.path.join(folder, file.replace(".pdf", ".pptx"))
                pdf_files.append((pdf_path, pptx_path, dpi))

    print(f"Found {len(pdf_files)} PDFs to convert.")

    for pdf_path, pptx_path, dpi in pdf_files:
        pdf_to_pptx(pdf_path, pptx_path, dpi)

    print("All conversions completed!")


convert_pdfs_in_directory("./directory", dpi=300)
